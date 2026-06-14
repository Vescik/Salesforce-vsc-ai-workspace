"""PPTX converter (requires ``python-pptx``).

Speaker notes are extracted to a sibling section per slide and marked with
``metadata['speaker_notes_present']=True`` so validators can enforce the
governance review described in the Knowledge 2.0 plan.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from ai_workspace.knowledge.converters import make_doc


def convert(path: Path | str) -> dict[str, Any]:
    source = Path(path)
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return make_doc(
            "pptx",
            source,
            warnings=["python-pptx not installed; .pptx skipped. Install salesforce-ai-workspace-tools[knowledge]."],
            parse_status="failed",
            degraded=True,
        )

    try:
        prs = Presentation(str(source))
    except Exception as exc:  # noqa: BLE001 - python-pptx raises a variety of types
        return make_doc("pptx", source, warnings=[f"python-pptx failed to open file: {exc}"], parse_status="failed")

    sections: list[dict[str, Any]] = []
    tables: list[dict[str, Any]] = []
    speaker_notes_present = False
    warnings: list[str] = []

    for index, slide in enumerate(prs.slides, start=1):
        anchor = f"slide-{index}"
        title_text, bullet_lines = _extract_slide_text(slide)
        heading = title_text or f"Slide {index}"
        body = "\n".join(bullet_lines).strip()
        sections.append({
            "heading": heading,
            "level": 2,
            "body": body,
            "anchors": [anchor],
        })

        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                rows: list[list[str]] = []
                for row in table.rows:
                    rows.append([cell.text.strip() for cell in row.cells])
                if rows:
                    tables.append({"caption": f"{heading} ({anchor})", "rows": rows})

        notes_text = _extract_notes(slide)
        if notes_text:
            speaker_notes_present = True
            sections.append({
                "heading": f"Slide {index} notes",
                "level": 3,
                "body": notes_text,
                "anchors": [f"{anchor}-notes"],
            })

    if speaker_notes_present:
        warnings.append(
            "Source PPTX contains speaker notes; mandatory governance review required before publish."
        )

    metadata: dict[str, Any] = {
        "library": "python-pptx",
        "slide_count": len(prs.slides),
        "speaker_notes_present": speaker_notes_present,
    }
    try:
        core = prs.core_properties
        metadata["author"] = str(core.author or "")
        metadata["title"] = str(core.title or "")
        if core.created:
            metadata["created"] = core.created.isoformat()
        if core.modified:
            metadata["modified"] = core.modified.isoformat()
    except Exception:  # noqa: BLE001
        pass

    return make_doc(
        "pptx",
        source,
        sections=sections,
        tables=tables,
        metadata=metadata,
        warnings=warnings,
        parse_status="ok",
        degraded=False,
    )


def _extract_slide_text(slide: Any) -> tuple[str, list[str]]:
    title_text = ""
    bullet_lines: list[str] = []
    if slide.shapes.title is not None:
        title_text = (slide.shapes.title.text or "").strip()
    for shape in slide.shapes:
        if shape == slide.shapes.title:
            continue
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            text = "".join(run.text for run in paragraph.runs).strip()
            if text:
                bullet_lines.append(text)
    return title_text, bullet_lines


def _extract_notes(slide: Any) -> str:
    try:
        if not slide.has_notes_slide:
            return ""
        notes_frame = slide.notes_slide.notes_text_frame
        if notes_frame is None:
            return ""
        text = (notes_frame.text or "").strip()
        return text
    except Exception:  # noqa: BLE001
        return ""
